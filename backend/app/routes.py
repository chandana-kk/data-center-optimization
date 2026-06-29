from io import BytesIO, StringIO
import csv
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response, StreamingResponse
from sqlalchemy.orm import Session
from app.core.security import create_access_token, hash_password, verify_password, get_current_user
from app.database import get_db
from app.dto import ChatRequest, LoginRequest, RegisterRequest, SandboxRequest, TokenResponse
from app.models import User
from app.repositories import AnalyticsRepository, UserRepository
from app.services.analytics import AnalyticsService

router = APIRouter(prefix="/api")


def service(db: Session = Depends(get_db)) -> AnalyticsService:
    return AnalyticsService(AnalyticsRepository(db))


@router.post("/auth/register", response_model=TokenResponse)
def register(payload: RegisterRequest, db: Session = Depends(get_db)):
    repo = UserRepository(db)
    if repo.get_by_email(payload.email):
        raise HTTPException(status_code=409, detail="Email already registered")
    user = repo.create(User(email=payload.email, name=payload.name, hashed_password=hash_password(payload.password)))
    return TokenResponse(access_token=create_access_token(user.email))


@router.post("/auth/login", response_model=TokenResponse)
def login(payload: LoginRequest, db: Session = Depends(get_db)):
    user = UserRepository(db).get_by_email(payload.email)
    if not user or not verify_password(payload.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return TokenResponse(access_token=create_access_token(user.email))


@router.post("/auth/forgot-password")
def forgot_password(payload: dict):
    return {"message": "Password reset instructions have been queued if the email exists."}


@router.get("/dashboard/summary")
def dashboard_summary(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.dashboard_summary()


@router.get("/infrastructure/analytics")
def infrastructure(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.infrastructure()


@router.get("/energy/analytics")
def energy(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.energy()


@router.get("/cooling/analytics")
def cooling(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.cooling()


@router.get("/carbon/analytics")
def carbon(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.carbon()


@router.get("/predictions/failures")
def predictions(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.predictions()


@router.get("/optimization/cloud")
def cloud_optimization(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.cloud_optimization()


@router.get("/recommendations")
def recommendations(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.recommendations()


@router.get("/digital-twin")
def digital_twin(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.digital_twin()


@router.post("/sandbox/simulate")
def sandbox(payload: SandboxRequest, svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.sandbox(payload)


@router.get("/executive")
def executive(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.executive()


@router.get("/extras")
def extras(svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.extras()


@router.post("/chat")
def chat(payload: ChatRequest, svc: AnalyticsService = Depends(service), current_user: str = Depends(get_current_user)):
    return svc.chat(payload.question)


@router.get("/reports/export/{format}")
def export_report(format: str, token: str | None = None, svc: AnalyticsService = Depends(service)):
    if token:
        try:
            get_current_user(token)
        except HTTPException:
            raise HTTPException(status_code=401, detail="Invalid token")
    summary = svc.dashboard_summary()
    if format == "csv":
        buffer = StringIO()
        writer = csv.writer(buffer)
        writer.writerow(["Metric", "Value"])
        for key, value in summary.items():
            writer.writerow([key, value])
        return Response(buffer.getvalue(), media_type="text/csv", headers={"Content-Disposition": "attachment; filename=neurodc-report.csv"})
    if format == "pdf":
        from reportlab.pdfgen import canvas
        buffer = BytesIO()
        pdf = canvas.Canvas(buffer)
        pdf.setTitle("NeuroDC Executive Report")
        pdf.drawString(72, 780, "NeuroDC Executive Optimization Report")
        y = 740
        for key, value in summary.items():
            pdf.drawString(72, y, f"{key}: {value}")
            y -= 22
        pdf.save()
        buffer.seek(0)
        return StreamingResponse(buffer, media_type="application/pdf", headers={"Content-Disposition": "attachment; filename=neurodc-report.pdf"})
    if format == "xlsx":
        import pandas as pd
        buffer = BytesIO()
        pd.DataFrame(summary.items(), columns=["Metric", "Value"]).to_excel(buffer, index=False)
        buffer.seek(0)
        return StreamingResponse(buffer, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers={"Content-Disposition": "attachment; filename=neurodc-report.xlsx"})
    if format == "pptx":
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "NeuroDC Executive Report"
        textbox = slide.shapes.add_textbox(914400, 1500000, 8000000, 4000000)
        frame = textbox.text_frame
        for key, value in summary.items():
            frame.add_paragraph().text = f"{key}: {value}"
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return StreamingResponse(buffer, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": "attachment; filename=neurodc-report.pptx"})
    raise HTTPException(status_code=400, detail="Supported formats: csv, pdf, xlsx, pptx")
